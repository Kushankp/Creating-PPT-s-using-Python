# -*- coding: utf-8 -*-
"""PPT.ipynb

Automatically generated by Colaboratory.

Original file is located at
    https://colab.research.google.com/drive/1GlzQfSqGyppaL_2D2vD5KrUvc1ngdeMc
"""

!pip install python-pptx

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Function to set font style
def set_font(text_frame, font_file, font_size):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.file = font_file
            run.font.size = Pt(font_size)

# Function to add content to a slide
def add_content_to_slide(slide, content, font_file, font_size):
    text_frame = slide.shapes[0].text_frame
    p = text_frame.add_paragraph()
    p.text = content
    p.alignment = PP_ALIGN.LEFT
    set_font(text_frame, font_file, font_size)

# Create a new presentation
presentation = Presentation()

# Slide 1
slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
add_content_to_slide(slide1, open("sample_slide1_input.txt").read(), "sample_font_file.ttf", 12)

# Slide 2
slide2 = presentation.slides.add_slide(presentation.slide_layouts[1])
add_content_to_slide(slide2, open("sample_slide2_input.txt").read(), "sample_font_file.ttf", 12)

# Save the presentation
presentation.save("output.pptx")

